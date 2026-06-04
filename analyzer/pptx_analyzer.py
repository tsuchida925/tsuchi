"""
PPTXテキスト抽出モジュール
"""

from typing import Tuple


def extract_pptx_text(file_bytes: bytes) -> Tuple[str, str]:
    """PPTXバイト列からテキストを抽出。(text, error)を返す"""
    try:
        from pptx import Presentation
        import io

        prs = Presentation(io.BytesIO(file_bytes))
        slides_text = []

        for i, slide in enumerate(prs.slides):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            parts.append(line)
                # テーブルのテキストも抽出
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            t = cell.text.strip()
                            if t:
                                parts.append(t)
            if parts:
                slides_text.append(f"[スライド{i+1}]\n" + "\n".join(parts))

        return "\n\n".join(slides_text), ""

    except ImportError:
        return "", "python-pptxが未インストール（pip install python-pptx）"
    except Exception as e:
        return "", f"PPTX解析エラー: {str(e)[:80]}"
